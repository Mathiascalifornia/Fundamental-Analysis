from pptx import Presentation
from pptx.util import Inches 
import os 

class PresPPT:

    def __init__(self):
        self.pres = Presentation('data\\template.pptx')

    def pres_title(self):
        '''Create a slide with the title'''
        layout = self.pres.slide_layouts[0] # Create layout
        slide = self.pres.slides.add_slide(layout) # Add layout 
        slide.shapes.title.text = self.title # Add the title


    def pres_description(self , string , title):
        '''Add the description of the company (On the first slide)'''
 
        layout = self.pres.slide_layouts[1]
        slide = self.pres.slides.add_slide(layout)
        slide.shapes.title.text = str(title)
        slide.shapes.add_textbox(Inches(0.00001) , Inches(1.75) , Inches(0.0001) , Inches(0.0001))\
                                .text_frame.add_paragraph().text = string
        


    def add_picture(self , picture_name , title,  left=1 , top=2):
        '''Function to add a picture on the presentation'''
        left_ =  Inches(left)
        top_ = Inches(top)
        layout = self.pres.slide_layouts[1] # Both title and content
        slide = self.pres.slides.add_slide(layout)
        slide.shapes.title.text = title
        pic = slide.shapes.add_picture(picture_name , left_ , top_)
        os.remove(picture_name)
