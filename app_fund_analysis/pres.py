from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

import os

path_template = os.path.join(os.path.dirname(os.path.dirname(__file__)), "data")


class PresPPT:

    data_path = os.path.join(os.path.dirname(os.path.dirname(__file__)), "data")

    def __init__(self):

        path_template = os.path.join(PresPPT.data_path, "template.pptx")
        self.pres = Presentation(path_template)

    def pres_title(self):
        """Create a slide with the title"""
        layout = self.pres.slide_layouts[0]  # Create layout
        slide = self.pres.slides.add_slide(layout)  # Add layout
        slide.shapes.title.text = self.title  # Add the title

    def pres_description(self, string, title):
        """Add the texts"""

        layout = self.pres.slide_layouts[1]
        slide = self.pres.slides.add_slide(layout)
        slide.shapes.title.text = str(title)

        textbox = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.75), Inches(8), Inches(1)
        )
        text_frame = textbox.text_frame
        text_frame.text = string

        title_shape = slide.shapes.title

        for paragraph in title_shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(25)

        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.LEFT

    # def add_picture(self , picture_name , title,  left=1 , top=2):
    #     '''Function to add a picture on the presentation'''
    #     left_ =  Inches(left)
    #     top_ = Inches(top)
    #     layout = self.pres.slide_layouts[1] # Both title and content
    #     slide = self.pres.slides.add_slide(layout)
    #     slide.shapes.title.text = title
    #     pic = slide.shapes.add_picture(picture_name , left_ , top_)
    #     os.remove(picture_name)

    def add_picture(self, picture_name, title, left=1, top=2):
        """
        Function to add a picture on the presentation and automatically center everything
        """

        # Get the current slide layout
        layout = self.pres.slide_layouts[1]  # Both title and content

        # Add a new slide with the current layout
        slide = self.pres.slides.add_slide(layout)

        # Add the image to the slide
        pic = slide.shapes.add_picture(picture_name, 0, 0)

        # Center the image horizontally
        pic.left = int((self.pres.slide_width - pic.width) / 2)

        # Center the image vertically
        pic.top = int((self.pres.slide_height - pic.height) / 2)

        # Define the title shape
        title_shape = slide.shapes.title

        # Set the title text
        title_shape.text_frame.text = title

        # Set the font size of the title
        for paragraph in title_shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(25)

        # Remove the image file after adding
        os.remove(picture_name)
