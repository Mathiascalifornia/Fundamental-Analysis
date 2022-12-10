############### Main function ###############
def run(language , ticker , path , companie_name):

    # Data manipulation
    import pandas as pd 
    import numpy as np
    import datetime as dt
    import textwrap
    import collections
    import collections.abc
    import talib as ta
    from googletrans import Translator
    import re 

    # Data visualisation
    import matplotlib.pyplot as plt
    import mplfinance as mpf
    import seaborn as sns 
    from pptx import Presentation
    from pptx.util import Inches 
    import dataframe_image as dsi
    

    # Scraping and API
    import requests
    from bs4 import BeautifulSoup
    from selenium import webdriver
    from selenium.webdriver.chrome.service import Service
    from selenium.webdriver.common.by import By 
    from selenium.webdriver.common.keys import Keys
    from webdriver_manager.chrome import ChromeDriverManager
    import yfinance as yf
    import pandas_datareader as web
    
    # Machine learning
    from keras.models import Sequential
    from keras.layers import Dense , BatchNormalization
    from sklearn.preprocessing import MinMaxScaler
    from vaderSentiment.vaderSentiment import SentimentIntensityAnalyzer
    import tensorflow as tf
    import keras

    # Other 
    import os
    import time
    import warnings


    # Filter the warnings 
    warnings.filterwarnings('ignore')
    
    # Translate in english
    t = lambda string : Translator().translate(string , src='French').text
    english = False
    if language == 'English':
        english = True



    ############### Get the URL , using the company's name ###############
    def get_url(ticker):
        ''' Get url using selenium. '''

        driver = webdriver.Chrome(service=Service(ChromeDriverManager().install()))#service=s
        #driver.get(f'https://www.google.com/search?q=zone+bourse+{companie_name}+{ticker}+finance&oq=zone+bourse+{ticker}+financz&aqs=chrome..69i57.7371j0j15&sourceid=chrome&ie=UTF-8')
        driver.get(f'https://www.google.com/search?q=zone+bourse+{companie_name}+finance&sxsrf=ALiCzsbIaWNWrnXJ5acLqlPx2kINT72YMA%3A1670610120483&ei=yHyTY9CSHcmPkdUP3veM2AI&oq=zone+bourse+telenor++finance&gs_lcp=Cgxnd3Mtd2l6LXNlcnAQAxgAMgQIIxAnOggIABCiBBCwA0oECEEYAUoECEYYAFCMBViMBWCHEGgBcAB4AIABKYgBKZIBATGYAQCgAQHIAQPAAQE&sclient=gws-wiz-serp')
        time.sleep(15) # Time to do the capcha if needed
        #driver.find_element(By.XPATH , ('//*[@id="L2AGLb"]')).click()
        try:
            driver.find_element(By.XPATH , ('//*[@id="L2AGLb"]')).click()
        except:
            pass
        try:
            driver.find_element(By.XPATH , ('//*[@id="rso"]/div[1]/div/div/div[1]')).click()
        
        except:
            try:
                driver.find_element(By.XPATH , ('//*[@id="rso"]/div[1]/div/div/div[1]/div')).click()
                
            except:
                driver.find_element(By.XPATH , ('//*[@id="rso"]/div[1]/div/div/div[1]/div/a')).click()
            
                

        current = driver.current_url
        if len(current) == 27:
            driver = webdriver.Chrome(service=Service(ChromeDriverManager().install()))#service=s
            driver.get(f'https://www.google.com/search?q=zone+bourse+{ticker}+finance&oq=zone+bourse+{ticker}+financz&aqs=chrome..69i57.7371j0j15&sourceid=chrome&ie=UTF-8')
            driver.find_element(By.XPATH , ('//*[@id="L2AGLb"]/div')).click()
            driver.find_element(By.XPATH , ('//*[@id="rso"]/div[2]/div/div/div[1]/div/a/h3')).click()
            current = driver.current_url

        if current[-1] == '?':
            current = current[0 : -1]

        if current.split('/')[-2] == 'consensus':
            current = '/'.join(current.split('/')[0 : -2]) + '/'

            
        return current


    

    # Get the url
    url_ =  str(get_url(ticker))
    print(url_)

    # Preprocess the urls
    if url_.split('/')[-2] != 'fondamentaux':
            url_desc = url_ + 'societe/'
    else:
            url_desc = url_.replace(url_.split('/')[-2] , 'societe')


    if url_.split('/')[-2] != 'fondamentaux':
            url = url_ + 'fondamentaux/'
    else:
            url = url_

    print(url)
    # Get the current year
    limit_year = int(int(dt.datetime.today().year) - 1)

    
    ############### Scrape description ###############

    def get_description(url=url_desc):

        '''Get the description of the company from the zonebourse site'''

        # Request
        responce =  requests.get(url)
        
        # Scrape raw data
        soup = BeautifulSoup(responce.text , 'html.parser')

        # Define a pattern with undesired element
        html_ = re.compile('<.*?>|&([a-z0-9]+|#[0-9]{1,6}|#x[0-9a-f]{1,6});\n')

        # Find the description
        soup_ = soup.find(class_='std_txt th_inner')

        # Clean 
        soup_ = re.sub(html_ , '' , str(soup_))

        if english:
            return t(soup_.replace('\n' , '').replace('\r' , ''))
        else:
            return soup_.replace('\n' , '').replace('\r' , '')


    


    # Jump line , for better fitting in the power point
    def jump_line(s : str , every=80):
        '''Jump line every 80 words , to better suits the power-point format'''
        s = '\n'.join(textwrap.wrap(s , every))
        if len(s) >= 960:
            s1 = s[:960]
            s2 = s[960:]
            return s1 , s2
        else:
            return s




    ############## Scrape titles , create the soup ###############
    responce = requests.get(url)
    soup = BeautifulSoup(responce.content , 'html.parser')
    list_table = ['Tableau_Histo_Valo' , 'Tableau_Histo_ECR_a' , 'Tableau_Histo_ECR_q' , 'Tableau_Histo_Sit_Fin']

    def get_titles(table , soup=soup):
        '''Get the titles (Capitalization , free cash flow ...) using the soup object'''
        to_del = re.compile('<td class="bc2T"|href=|"/formation/espace_pedagogique/La-terminologie-et-les-ratios-boursiers-133/">Capitalisation</a><sup|"/formation/espace_pedagogique/La-terminologie-et-les-ratios-boursiers-133/"|"/formation/espace_pedagogique/La-terminologie-comptable-132/"|class=|</td>|</a>|><a|>|title=|</sup|<td|"bc2T bc2gt"|<sup|<i|</i|"USD en Millions"|"USD"|1|2|3|CAD en Millions|EUR en Millions"|"|CHF|en Millions|en millions|INR|GBP|NOK|CNY|Sek|SEK')
        titles = soup.find(id=table).find_all(class_='bc2T')

        title = []
        for value in titles:
            title.append(re.sub(to_del , '' , str(value)))

        if str(title[0]).replace(' ' , '') == '':    
            title[0] = 'Capitalisation'

        if english:
            title = [t(tit) for tit in title]   

        return title



    ############### Scrape the numbers ##############
    def get_numbers(table , soup=soup):
        '''Get the numbers related to the titles'''
        # To delete
        to_del = re.compile('" title="Nombre d\'analystes financiers ayant fourni une estimation : (\\d\\d|\\d)|class="bc2V tableCol(\\d\\d|\\d) bc2H"|class="bc2V tableCol(\\d\\d|\\d)"|class="bc2V tableCol(\\d\\d|\\d) bc2g"|class="bc2V tableCol(\\d\\d|\\d) bc2g bc2H"|style=|</td>|<td|Dernière mise à jour')

        # Define a pattern
        pattern = re.compile('(bc2V tableCol\\d|bc2V tableCol\\d bc2H|bc2V tableCol\\d bc2g)')

        # Find the numbers
        number = soup.find(id=table).find_all(class_=pattern)

        # Clean and append the numbers
        numbers=[]
        for n in number:
            numbers.append(re.sub(to_del , '' , str(n)))

        # Second cleaning
        to_del = re.compile('  "background-color:#DEFEFE;">|"background-color:#DEFEFE;display:none;"|"display:none;"|title=| |"|>|-|:')
        clean_n0 = []
        for n in numbers:
            clean_n0.append(re.sub(to_del , '' , str(n)))


        # Replace the empty values by np.nan , append the data
        clean_n = []
        for val in clean_n0:
            if len(str(val).lstrip()) == 0:
                clean_n.append(np.nan)
            else:
                clean_n.append(str(val).replace(',' , '.'))



        # Create the one date pattern
        pattern = re.compile('^\\d\\d/\\d\\d/\\d\\d\\d\\d$')

        # Replace the bad dates by np.nan
        for i in range(len(clean_n)):
            if re.findall(pattern , str(clean_n[i])) != []:
                clean_n[i] = np.nan
                
            
        # Create the two date pattern
        pattern = re.compile('^\\d\\d/\\d\\d/\\d\\d\\d\\d\\d\\d/\\d\\d/\\d\\d\\d\\d$')

        # Kill the replicates
        for i in range(len(clean_n)):
            if re.findall(pattern , str(clean_n[i])) != []:
                clean_n[i] = clean_n[i][0 : 10]

        return clean_n



    ########## Years ##########

    def get_years(table , soup=soup):
        '''Get the indicatprs's years'''
        # Define pattern
        pattern = re.compile('bc2Y tableCol\\d')

        # Get the years 
        dates = soup.find(id=table).find_all(class_=pattern)
        
        # Create a pattern for years
        year = re.compile('(\\d\\d\\d\\d\\s[A-Z]\\d|\\d\\d\\d\\d)')

        # Append them
        years = []
        for date in dates:
            years.append(re.findall(year , str(date)))

        # Clean the years column
        clean_year = []
        for year in years:
            clean_year.append(str(year).replace('[' , '').replace(']' , '').replace("'" , ""))

        return clean_year

 
    ########## Key numbers dataframe ##########

    def get_key_df(soup=soup):
        '''Get the key number dataframe on the zone bourse site'''

        # Get the table
        table = soup.find(attrs={'class':'Bord'})

        # Create the delete pattern
        pattern = re.compile('<tr>|<table|<td>|</td>|</tr>|style="text-align:|\n|"|<td|right')


        # Append the raw values
        raw_values = []
        for val in table:
            values = re.sub(pattern , '' , str(val))
            if len(values) > 1:
                raw_values.append(values)

        # Define patterns to catch everything before '>'
        pattern_b = re.compile('(.*?)>')
        pattern_a = re.compile('\>(.*)')
        key_titles = []
        key_numbers = []

        # Get the separate values
        for val in raw_values:
            key_titles.append(re.findall(pattern_b , str(val)))
            key_numbers.append(re.findall(pattern_a , str(val)))

        # Verify that everything is in order
        assert len(key_titles) == len(key_numbers)

        # Assemble the dataframe
        if english:
            df_key = pd.DataFrame({'Indicator' : key_titles , 'Number' : key_numbers})
            df_key['Indicator'] = df_key['Indicator'].apply(lambda x : t(str(x).replace('[' , '').replace(']' , '').replace("'" , '').replace('"' , '')))
            df_key['Number'] = df_key['Number'].apply(lambda x : t((str(x).replace('[' , '').replace(' ' , '').replace(']' , '').replace('%' , '').replace(',' , '.').replace("'" , ''))))

        else:
            df_key = pd.DataFrame({'Indicateur' : key_titles , 'Nombre' : key_numbers})
            df_key['Indicateur'] = df_key['Indicateur'].apply(lambda x : str(x).replace('[' , '').replace(']' , '').replace("'" , '').replace('"' , ''))
            df_key['Nombre'] = df_key['Nombre'].apply(lambda x : (str(x).replace('[' , '').replace(' ' , '').replace(']' , '').replace('%' , '').replace(',' , '.').replace("'" , '')))
        
        return df_key


    ############### Assemble the dataframe elements ###############

    def assemble_dataframe(title , values , column):   
        '''Assemble the dataframe , using the titles , numbers and years'''

        # Lenght of the column , to slice the values
        n = len(column)

        # Preprocess the values , make list of lists
        ind = [0]
        #num = n

        # Create a list of indices
        for i in range(len(title)):
            ind.append(n + i*n)
        
        # Slice the values with the new list of indices
        list_val = []
        for i in range(len(ind)):
            try:

                list_val.append(values[ind[i] : ind[i+1]])
            except IndexError:
                pass 

        # Create the dataframe
        df = pd.DataFrame(index=title , data=list_val , columns=column)

        # Drop the publication dates
        df.drop('Date de publication' , axis=0 , inplace=True , errors='ignore')
        df.drop('Publication date' , axis=0 , inplace=True , errors='ignore')

        return df

    ############### Get stock data ###############

    # Stock price 
    def get_price(tickers  ,  start : dt.datetime , data_source ='yahoo'):
        '''Get the stock price , volume , open , close , high , low ... Using Datareader (data from yahoo finance)'''
        return web.DataReader(tickers , data_source ,  start)

    # Dividends
    def get_dividend(tickers : str , start : dt.datetime , data_source='yahoo-dividends'):
        '''Get the dividend history'''
        return web.DataReader(tickers , data_source ,  start)['value']


    # Mean return 
    def annualized_return(df : pd.DataFrame):
        '''Get the mean annualized return'''
        df_ = df.copy()
        df_ = df_['Adj Close']
        years = pd.to_datetime(dt.date.today()).year - df_.index[0].year

        # Total returns
        total_return = (df_.iloc[-1] - df_.iloc[0]) / df_.iloc[0]

        # Annualized return
        return ((total_return + 1)**(1/years)) -1 , years


    # Annualized mean return for the last 5 years
    def annualized_return_five_years(df : pd.DataFrame):
        '''Get the mean annualized return for the last five years'''
        df_ = df.copy()
        df_ = df_['Adj Close']
        new_index = df_.index[-1] - dt.timedelta(252*5)
        df_ = df_[new_index:]
        years = pd.to_datetime(dt.date.today()).year - df_.index[0].year

        # Total returns
        total_return = (df_.iloc[-1] - df_.iloc[0]) / df_.iloc[0]

        # Annualized return
        return ((total_return + 1)**(1/years)) -1 , years


    # Sharpe and sortino ratio
    def sharpe_and_sortino_ratio(df : pd.DataFrame , rfr=0 , target=0):
        '''Return the sharpe and sortino ratio'''
        df_ = df.copy()
        df_ = df_['Adj Close']
        # Number of years
        years = pd.to_datetime(dt.date.today()).year - df_.index[0].year

        # Sharpe ratio
        total_return = (df_.iloc[-1] - df_.iloc[0]) / df_.iloc[0]
        annualized_return = ((1 + total_return) ** (1/years)) -1
        returns = df_.pct_change()
        vol = returns.std() * np.sqrt(252)
        sharpe = ((annualized_return - rfr) / vol)

        # Sortino ratio
        downside_return = returns[returns.values < target]
        mean_return= returns.mean()
        vol_down = downside_return.std()
        sortino = (mean_return-rfr) / vol_down

        return sharpe , sortino



    def sharpe_and_sortino_ratio_five_years(df : pd.DataFrame , rfr=0 , target=0):
            '''Returns the sharpe and sortino ratio for the last 5 years'''
            df_ = df.copy()
            df_ = df_['Adj Close']
            new_index = df_.index[-1] - dt.timedelta(252*5)
            df_ = df_[new_index:]

            # Number of years
            years = pd.to_datetime(dt.date.today()).year - df_.index[0].year

            # Sharpe ratio
            total_return = (df_.iloc[-1] - df_.iloc[0]) / df_.iloc[0]
            annualized_return = ((1 + total_return) ** (1/years)) -1
            returns = df_.pct_change()
            vol = returns.std() * np.sqrt(252)
            sharpe = ((annualized_return - rfr) / vol)

            # Sortino ratio
            downside_return = returns[returns.values < target]
            mean_return= returns.mean()
            vol_down = downside_return.std()
            sortino = (mean_return-rfr) / vol_down

            return sharpe , sortino



    # Main investors
    def main_institutions(ticker=ticker):
        '''Returns the three biggest institutional holders'''
        data_ = yf.Ticker(str(ticker).upper())
        df__ = data_.institutional_holders
        return list(df__['Holder'][0 : 3])




    ############### Create a power point report #############

    # Define an url for saving the picture
    url_pic = '\\'.join(path.split('\\')).replace('"' , '') #"C:\\Users\\lucas\\OneDrive\\Bureau\\Code\\Python\\DATACAMP\\Projets\\analyse_entreprise\\Presentation"

    # Create a presentation object
    pres = Presentation('data\\template.pptx')

    # Create a slide with the title
    def title(title , pres=pres):
        # Create layout
        layout = pres.slide_layouts[0]
        # Add layout 
        slide = pres.slides.add_slide(layout)
        # Add the title
        slide.shapes.title.text = title


    def description(string, title, pres=pres):
        '''Add the description of the company (On the first slide)'''
        # Description layer
        layout = pres.slide_layouts[1]
        slide = pres.slides.add_slide(layout)
        slide.shapes.title.text = str(title)
        slide.shapes.add_textbox(Inches(0.00001) , Inches(1.75) , Inches(0.0001) , Inches(0.0001)).text_frame.add_paragraph().text = string


    def add_picture(picture_name , title , pres=pres , left=1 , top=2):
        '''Function to add a picture on the presentation'''
        left_ =  Inches(left)
        top_ = Inches(top)
        layout = pres.slide_layouts[1] # Both title and content
        slide = pres.slides.add_slide(layout)
        slide.shapes.title.text = title
        pic = slide.shapes.add_picture(picture_name , left_ , top_)
        os.remove(picture_name)



    ############## Visualisation ###############

    # Plot df elements
    def plot_element(df : pd.DataFrame):
        '''Function to plot the elements of the indicators dataframes over time'''
        cleaning = lambda x : str(x).replace('x' , '').replace('%' , '')
        for col in df.columns:
            new_col = []
            for val in df[col].apply(cleaning):
                    new_col.append(float(val))

            df[col] = new_col
        
        with sns.plotting_context('talk'):

            sns.set_style('darkgrid')
            index = list(df.index)
            for i in range(len(index)):
                plt.figure(figsize=(8,5))
                d = df.iloc[i]
                d = d.fillna(method='pad')
                d.index = pd.to_datetime(d.index).year
                preds = d[d.index > limit_year]
                plt.plot(d , color='red')
                plt.plot(preds , linestyle='--' , color='white' , label='Predictions')
                plt.legend()
                plt.savefig('data\\fig{i}.png')
                plt.close('all')
                add_picture('data\\fig{i}.png' , index[i] , left=1.1 , top=1.9)
                

            
    # Linear regression , with stock data
    def plot_regression(df : pd.DataFrame):
        '''Plot the stock price with a linear regression , to see the trend'''
        
        df_ = df.copy()
        df_.drop(['High', 'Low', 'Open', 'Close', 'Volume'] , axis=1 , inplace=True)
        df_['day_from_start'] = (df_.index - df_.index[-1]).days
        X = df_['day_from_start'].values
        y = df_.drop('day_from_start' , axis=1).values
        a , b = np.polyfit(X , y , deg=1)
        x_line = np.array([np.min(X) , np.max(X)]) # Just the lenght of the regression line
        y_line = a*x_line + b

        with sns.plotting_context('notebook'):
            
            sns.set_style('darkgrid')
            plt.figure(figsize=(8,5))
            plt.plot(X , y ,color='blue')
            plt.plot(x_line , y_line , color='red')
            plt.savefig('data\\linear_regression.png')
            plt.close('all')
            if english:
                add_picture('data\\linear_regression.png' , 'Linear regression')
            else:
                add_picture('data\\linear_regression.png' , 'Régression linéaire')
            


    # Plot maximum draw down
    def plot_maximum_draw_down(df : pd.DataFrame):
        ''' Plot the maximum drow down , which is a measure of an asset's largest price drop from a peak to a trough'''

        _df_ = df.copy()
        _df_.drop(['High', 'Low', 'Open', 'Close', 'Volume'] , axis=1 , inplace=True , errors='ignore')
        # Calculate the max value 
        roll_max = _df_.rolling(center=False,min_periods=1,window=252).max()

        # Calculate the daily draw-down relative to the max
        daily_draw_down = _df_/roll_max - 1.0

        # Calculate the minimum (negative) daily draw-down
        max_daily_draw_down = daily_draw_down.rolling(center=False,min_periods=1, window=252).min()

        # Plot the results
        with sns.plotting_context('notebook'):

            sns.set_style('darkgrid')
            plt.figure(figsize=(8,5))
            plt.plot(daily_draw_down.index, daily_draw_down, label='Daily drawdown')
            plt.plot(max_daily_draw_down.index, max_daily_draw_down, label='Maximum daily drawdown in time-window' , color='red')
            plt.savefig('data\\maximum_draw_down.png')
            plt.legend()
            plt.close('all')
            if english:
                add_picture('data\\maximum_draw_down.png' , 'Maximum loss')
            else:
                add_picture('data\\maximum_draw_down.png' , 'Perte de valeur maximum')
            



    # Plot the stock data with the dividend
    def price_with_dividends(stock_price_ , div):
        '''Plot the stock price with the dividends payed over time'''
        try:
            _div = div.to_frame()
        except:
            pass

        df_ = stock_price_.copy()
        df_.drop(['High', 'Low', 'Open', 'Close', 'Volume'] , axis=1 , inplace=True , errors='ignore')
        with sns.plotting_context('notebook'):
            sns.set_style('dark')
            fig , ax = plt.subplots(figsize=(8,5))
            ax.plot(df_.index , df_.values , label='Stock price')
            plt.legend(loc='upper left')
            ax.twinx().plot(_div.index , _div.values , color='red' , label='Dividend')
            plt.legend(loc='center left')
            plt.savefig('data\\price_with_dividends.png')
            plt.close('all')
            if english:
                add_picture('data\\price_with_dividends.png' , 'Price and dividends')
            else:
                add_picture('data\\price_with_dividends.png' , 'Prix et dividendes')




    # Annual dividend per year
    def annual_dividend_history(div_):
        '''Group the dividend per year to see the annual dividend history'''
        try:
            div_ = div_.to_frame()
        except:
            pass

        div_['year'] = div_.index.year
        d__ = div_.groupby('year').sum()
        d__['year'] = d__.index

        with sns.plotting_context('notebook'):
            sns.set_style('darkgrid')
            plt.figure(figsize=(8,5))
            sns.barplot(x='year' , y='value' , data=d__ , edgecolor='black')
            plt.ylabel('Montant')
            plt.xlabel('Année')
            plt.xticks(rotation = 90)
            plt.savefig('data\\annual_dividend.png')
            plt.close('all')
            if english:
                add_picture('data\\annual_dividend.png' , 'Dividend per year and per share')
            else:
                add_picture('data\\annual_dividend.png' , 'Dividende versé par année et par action')


    # Pay out ratio
    def payout_ratio(div : pd.DataFrame , df2 : pd.DataFrame):   
        '''Function to compute and plot the payout ratio'''

        try :
            div = div.to_frame()
        except:
            pass
        if english == False:
            div['year'] = list(div.index.year.astype(str))
            d_ = div.groupby('year').sum()
            d_['year'] = d_.index
            years = [val for val in df2 if val in d_.index]
            df2.index = [str(val).rstrip().lstrip() for val in df2.index]
            bna_years = dict(df2[years].loc['BNA'])
            d_ = d_[d_['year'].isin(years)].drop('year' , axis=1 , errors='ignore')
            d_['BNA'] = list(bna_years.values())
            d_['payout'] = d_['value'] / d_['BNA'] * 100
        else:
            div['year'] = list(div.index.year.astype(str))
            d_ = div.groupby('year').sum()
            d_['year'] = d_.index
            years = [val for val in df2 if val in d_.index]
            df2.index = [str(val).rstrip().lstrip() for val in df2.index]
            bna_years = dict(df2[years].loc['Bna'])
            d_ = d_[d_['year'].isin(years)].drop('year' , axis=1 , errors='ignore')
            d_['Bna'] = list(bna_years.values())
            d_['payout'] = d_['value'] / d_['Bna'] * 100

        with sns.plotting_context('talk'):
            sns.set_style('darkgrid')
            plt.figure(figsize=(8,5))
            plt.plot(d_.index , d_['payout'] , color='red')
            plt.savefig('data\\payout.png')
            plt.close('all')
            if english:
                add_picture('data\\payout.png' , 'Pay out ratio')
            else:
                add_picture('data\\payout.png' , 'Taux de distribution')



# Plot against the SP500 returns
    def plot_against_benmark(df : pd.DataFrame , bench):
        '''Plot the normalized data against a benchmark (the sp500)'''
        df_ = df.copy()
        bench_ = bench.copy()
        df_['normalized'] = MinMaxScaler().fit_transform(df_['Adj Close'].values.reshape(-1,1))   
        bench_['normalized'] = MinMaxScaler().fit_transform(bench_['Adj Close'].values.reshape(-1,1))   

        with sns.plotting_context('notebook'):
            sns.set_style('darkgrid')
            plt.figure(figsize=(8,5))
            plt.plot(df_.index , df_['normalized'] , color='blue' , label=ticker)
            plt.plot(bench_.index , bench_['normalized'] , label='SP500' , color='red')
            plt.legend()
            plt.savefig('data\\bench.png')
            plt.close('all')
            if english:
                add_picture('data\\bench.png' , f"Normalized stocks price :{ticker} vs SP500")
            else:
                add_picture('data\\bench.png' , f"Prix de l'action normalisé {ticker} vs SP500")


    # Plot the propers fonds versus the debt for the five last years , plus prediction for the next years
    def cap_vs_debt(df : pd.DataFrame):
        '''Plot the proper capital versus the debt. In a perfect world , both should be more on less equivalent'''

        __df = df.copy()
        __df.index = [str(val).lstrip().rstrip() for val in df.index]
        if english == False:
            dette = __df.loc['Dette Nette']
            capitaux = __df.loc['Capitaux Propres']
        else:
            dette = __df.loc['Net debt']
            capitaux = __df.loc['Equity']

        debt = pd.DataFrame({'year' : dette.index[1:] , 'debt' : dette.values[1:]})
        cap = pd.DataFrame({'year' : capitaux.index[1:] , 'cap' : capitaux.values[1:]})
        new_df = debt.merge(cap ,on='year', how='inner')

        if new_df['debt'].isna().sum() <= 2 and new_df['cap'].isna().sum() <=2:

            new_df['debt'].fillna(method='pad' , inplace=True)
            new_df['cap'].fillna(method='pad' , inplace=True)
            new_df['debt'] = new_df['debt'].apply(lambda x : float(x))
            new_df['cap'] = new_df['cap'].apply(lambda x : float(x))
            new_df['year'] = new_df['year'].apply(lambda x : int(x))
            prediction = new_df[new_df['year'] > limit_year]
            
            with sns.plotting_context('talk'):

                sns.set_style('darkgrid')
                plt.figure(figsize=(8,5))
                if english == False:
                    plt.plot(new_df['year'] , new_df['debt'] , color='red' , label='Dette')
                    plt.plot(new_df['year'] , new_df['cap'] , color='blue' , label='Capitaux propres')
                    plt.plot(prediction['year'] , prediction['debt'] , linestyle='--' , color='white' , label='prediction')
                    plt.plot(prediction['year'] , prediction['cap'] , linestyle='--' , color='white')
                    plt.legend()
                    plt.savefig('data\\cap_versus_debt.png')
                    plt.close('all')
                    add_picture('data\\cap_versus_debt.png' , 'Capitaux propres vs Dette net')
                else:
                    plt.plot(new_df['year'] , new_df['debt'] , color='red' , label='Debt')
                    plt.plot(new_df['year'] , new_df['cap'] , color='blue' , label='Equity')
                    plt.plot(prediction['year'] , prediction['debt'] , linestyle='--' , color='white' , label='prediction')
                    plt.plot(prediction['year'] , prediction['cap'] , linestyle='--' , color='white')
                    plt.legend()
                    plt.savefig('data\\cap_versus_debt.png')
                    plt.close('all')
                    add_picture('data\\cap_versus_debt.png' , 'Equity versus debt')



    # Plot the shareholder pie
    def shareholders(ticker=ticker):
        '''Plot a pie plot of the three biggest institutions shareholders'''
        data_ = yf.Ticker(str(ticker).upper())
        d_f = data_.major_holders
        share_holders = d_f.iloc[0 : 2][0].apply(lambda x : float(str(x).replace('%' , '')))
        particuliers = 100 - share_holders.sum()
        color = ['red' , 'lightgreen' , 'gold']
        explode = [0.12 , 0.02 , 0.05]
        if english == False:
            labels = [f'Initiés : {share_holders[0]}%' , f'Institutions:{share_holders[1]} %' , f'Autre : {np.round(particuliers , 3)}%']
        else:
            labels = [f'Initiates : {share_holders[0]}%' , f'Institutions:{share_holders[1]} %' , f'Other : {np.round(particuliers , 3)}%']


        plt.figure(figsize=(5,5))
        plt.pie([share_holders[0] , share_holders[1] , particuliers] , colors=color , explode=explode)
        plt.legend(labels)
        plt.savefig('data\\shareholders.png')
        plt.close('all')
        if english == False:
            add_picture('data\\shareholders.png' , 'Repartition des investisseurs' , left=2.4 , top=1.5)
        else:
            add_picture('data\\shareholders.png' , 'Investor distribution' , left=2.4 , top=1.5)


    # RSI
    def plot_rsi(df : pd.DataFrame):
            '''Plot the Relative Strenght Index'''
            start = dt.datetime.now() - dt.timedelta(800)
            df['RSI'] = ta.RSI(df['Adj Close'])
            plt.figure(figsize=(10,3))
            plt.plot(df['RSI'][start:])
            plt.axhline(y=70 , color='red' , linestyle='--')
            plt.axhline(y=30 , color='green' , linestyle='--')
            plt.savefig('data\\rsi.png')
            plt.close('all')
            add_picture('data\\rsi.png' , 'Relative Strength Index (RSI)' , left=0 , top=2.5)


    # Zoom on the last six months
    def plot_zoom_candles(df : pd.DataFrame):
            '''Plot a zoom of the last six month , with candles and volume over time'''
            start = dt.datetime.now() - dt.timedelta(175)
            df['CDL'] = ta.CDLENGULFING(df['Open'] , df['High'] , df['Low'] , df['Close'])
            colors = mpf.make_marketcolors(up='green' , down='red')
            style = mpf.make_mpf_style(base_mpf_style='yahoo' , marketcolors=colors)
            ax = mpf.plot(df[start:] , type='candle' , style=style , figsize=(9,6) , volume=True , savefig='data\\Zoom.png')
            if english == False:
                add_picture('data\\Zoom.png' , 'Zoom des six derniers mois' , left=0.05 , top=1.5)
            else:
                add_picture('data\\Zoom.png' , 'Last six months zoom' , left=0.05 , top=1.5)


    ############### Neural network ###############

    def NN_model(df_ : pd.DataFrame , sp500 : pd.DataFrame , vix : pd.DataFrame):   
        '''Deep learning model to give a estimation of the direction the price might take for the next 5 years. Plot a visualisation of the results.'''

        if len(df_) >= 1260*3:
            
            # Get and preprocess the data
            int_rates = pd.read_csv("data\\FEDFUNDS.csv" , parse_dates=['DATE'] , index_col=['DATE'])
            gpd = pd.read_csv("data\\API_NY.GDP.MKTP.CD_DS2_en_csv_v2_4578252.csv")
            gpd = gpd[gpd['Country Code'] == 'USA'].drop(['Country Name' , 'Country Code' , 'Indicator Code' , 'Unnamed: 66' , 'Indicator Name'] , axis=1 , errors='ignore')
            inflation = pd.read_csv("data\\API_FP.CPI.TOTL.ZG_DS2_en_csv_v2_4578133.csv")
            inflation = inflation[inflation['Country Code'] == 'USA'].drop(['Country Name' , 'Country Code' , 'Indicator Name' , 'Indicator Code' , 'Unnamed: 66'] , axis=1 , errors='ignore')

            gpb_val = []
            for l in gpd.values:
                for j in l:
                    gpb_val.append(j)

            infl_val = []
            for l in inflation.values:
                for j in l:
                    infl_val.append(j)

            inf_col = list(inflation.columns)
            gpb_col = list(gpd.columns)
            gpb_col.append('2022') # To update every year 
            inf_col.append('2022') # To update every year
            infl_val.append(7.7) # To update every year
            gpb_val.append(22996100000000.0) # To update every year

            inflation = pd.DataFrame(index=inf_col , data=infl_val)
            gpd = pd.DataFrame(index=gpb_col , data=gpb_val)


            sp500 = sp500[df_.index[0] : df_.index[-1]]
            df_ = df_[vix.index[0]:]
            vix = vix[df_.index[0] : df_.index[-1]]

            gpd.index = pd.to_datetime(gpd.index)
            gpd = gpd.resample('B').interpolate()
            gpd = gpd[df_.index[0] : df_.index[-1]]
            inflation.index = pd.to_datetime(inflation.index)
            inflation = inflation.resample('B').interpolate()
            inflation = inflation[df_.index[0] : df_.index[-1]]

            # Merge the dataframes
            df = df_.merge(sp500 , left_index=True , right_index=True , how='left').merge(vix , left_index=True , right_index=True , how='left').merge(inflation , left_index=True , right_index=True , how='left').merge(gpd , left_index=True , right_index=True , how='left')
            df.fillna(method='pad' , inplace=True)
            df.columns = ['vol_t' , 'close_t' , 'vol_sp' , 'close_sp' , 'close_v' , 'inflation' , 'gpd']


            # Normalized data 
            scaled_df = MinMaxScaler().fit_transform(df)

            target = []
            s_train = []
            slice_ = 1260 # Five years in business days


            for i in range(1 , len(df.index) - slice_):

                stop = slice_ * i
                start = stop - slice_

                target.append(df['close_t'][start + slice_ : stop + slice_].values)
                s_train.append(scaled_df[start : stop])


            s_train = np.array([val for val in s_train if len(val) > 0])
            target = np.array([val for val in target if len(val) > 0])

            sresid_train = s_train[-1]
            s_train = s_train[0 : -1]
            x_future = s_train[-1]
            s_train = s_train[0 : -1]
            target = target[0 : -1]


            # Loss function
            def penalty_loss(y_true , y_pred):

                    penalty = 150.
                    loss = tf.where(tf.less(y_pred - y_true , 0), # If ...
                                            penalty * tf.square(y_true - y_pred), # Then ...
                                            tf.square(y_true - y_pred)) # Else ...
                            
                    return tf.reduce_mean(loss , axis=-1)

            # Add the custom loss to keras
            keras.losses.penalty_loss = penalty_loss


            # Model building
            model = Sequential()
            model.add(Dense(500 , input_shape=(7,) , activation='relu'))
            model.add(BatchNormalization())
            model.add(Dense(500 , activation='relu'))
            model.add(BatchNormalization())
            model.add(Dense(250 , activation='relu'))
            model.add(BatchNormalization())
            model.add(Dense(125 , activation='relu'))
            model.add(BatchNormalization())
            model.add(Dense(75 , activation='relu'))
            model.add(BatchNormalization())
            model.add(Dense(1 , activation='linear'))
            model.compile(optimizer='adam' , loss=penalty_loss)


            # Training loop
            for i in range(len(target)):
                model.fit(s_train[i] , target[i] , epochs=35)


            # Prediction
            preds = model.predict(x_future)
            replacement_val = df[dt.datetime.today() - dt.timedelta(1260):]['close_t'].mean()
            preds = [float(replacement_val / 2) if val <= (replacement_val / 2.5) else float(val) for val in preds]
            preds = [float(replacement_val * 1.5) if val > (replacement_val*2) else float(val) for val in preds]
            new_index = pd.date_range(start=df.index[-1] , periods=len(preds) , freq='B')
            to_plot_df = df['close_t']
            preds_df = pd.DataFrame(index=new_index , data=preds)
            to_plot = pd.concat([to_plot_df , preds_df])


            # Visualisation
            with sns.plotting_context('notebook'):
                sns.set_style('darkgrid')
                plt.figure(figsize=(8,5))
                plt.plot(to_plot)
                plt.plot(preds_df , color='red')
                if english == False:
                    plt.title('Estimation (en rouge) sur cinq ans , par intelligence artificielle.')
                else:
                    plt.title('AI estimate in red , for five years.')
                plt.savefig('data\\Prédictions.png')
                plt.close('all')
                if english == False:
                    add_picture('data\\Prédictions.png' , 'Estimation par IA')
                else:
                    add_picture('data\\Prédictions.png' , 'AI estimate')


    ############### Sentiment analysis using vader ###############

    def sentiment_scores(ticker=ticker):
        '''Sentiment analysis using vader SentimentAnalyser (code source modified to fit my needs) , using finviz news title concerning the companie.'''
        if '.' in ticker:
            ticker_ = ticker.split('.')[0]
            url = f'https://finviz.com/quote.ashx?t={ticker_}&p=d'
        else:
            url = f'https://finviz.com/quote.ashx?t={ticker}&p=d'

        driver = webdriver.Chrome(service=Service(ChromeDriverManager().install()))
        driver.get(url)
        time.sleep(5)
        driver.find_element(By.XPATH , ('//*[@id="qc-cmp2-ui"]/div[2]/div/button[3]')).click()
        time.sleep(3)
        source = driver.page_source
        driver.quit()


        source = BeautifulSoup(source)
        news = source.find_all(class_="news-link-container")

        pattern = re.compile('.*?\"_blank">(.*)</a>.*')

        data = []
        for new in news:
            data.append(str(re.findall(pattern , str(new))).replace('[' , '').replace(']' , '').replace("'" , '').replace('"' , ""))

        new_words = {
        'crushes': 10,
        'beats': 10,
        'misses': -10,
        'trouble': -10,
        'falls': -10,
        'down' : -15,
        'up' : 15,
        'go down' : -15,
        'go up' : 15,
        'increased' : 10,
        'decreased' : -10,
        'increase' : 10,
        'decrease' : -10,
        'safe' : 15,
        'unsafe' : -15,
        'attractive' : 10,
        'increase in value' : 10,
        'decrease in value' : -10,
        'top' : 15,
        'buy' : 5,
        'raise' : 5,
        'low' : -5,
        'dangerous' : -10 }

        vader = SentimentIntensityAnalyzer()
        vader.lexicon.update(new_words)

        to_plot = []

        for news in data:
            score = vader.polarity_scores(news)
            score.pop('compound' , [])
            if score['pos'] > score['neg']:
                to_plot.append('Positive headline')
                
            elif score['pos'] < score['neg']:
                to_plot.append('Negative headline')

        scores = pd.DataFrame(data=to_plot , columns=['score'])

        with sns.plotting_context('talk'):
            sns.set_style('darkgrid')
            plt.figure(figsize=(8,5))
            sns.countplot(x='score' , data=scores , palette={'Positive headline' : 'green' , 'Negative headline' : 'red'} , edgecolor='black')
            plt.xlabel('')
            plt.savefig('data\\sentiment.png')
            plt.close('all')
            if english == False:
                add_picture('data\\sentiment.png' , 'Analyse de sentiment du marché, basée sur les titres d\'actualités.')
            else:
                add_picture('data\\sentiment.png' , t('Analyse de sentiment du marché , basée sur les titres d\'actualités.'))



    ############### Create the presentation ###############
    try:

        if english == False:
            string = get_description()
            string += '\nLa capitalisation , valeur entreprise , chiffre d\'affaire , EBITDA , EBIT , EBT , le résultat net , la dette et trésorerie net , le free cash flow , les capitaux propres , le total des actifs et le Capex sont en millions.'
        else:
            string = get_description()
            string += t('\nLa capitalisation , valeur entreprise , chiffre d\'affaire , EBITDA , EBIT , EBT , le résultat net , la dette et trésorerie net , le free cash flow , les capitaux propres , le total des actifs et le Capex sont en millions.')

        if len(string) < 960:
            string = jump_line(string)
            description(string , title=str(url.split('/')[-3]))
        else:
            string1 , string2 = jump_line(string)
            description(string1 , title=str(url.split('/')[-3]))
            description(string2 , title='Suite :')

    except:

        pass
    
    
    worked_stock_price = False
    worked_dividends = False
    worked_dataframe = False
    worked_share_holders = False
    
    try:
        # Get the prices 
        stock_price = get_price(ticker , start=dt.datetime(1975,1,2))
        # SP500
        sp500 = get_price(['^GSPC'] , start=stock_price.index[0])
        # Vix
        vix = web.DataReader('^VIX' , 'yahoo' , stock_price.index[0])['Adj Close'].to_frame()

        worked_stock_price = True

    except:

        pass

    try:
        # Get the dividends
        dividends = get_dividend(ticker , start=dt.datetime(1990,1,2))
        if len(dividends) > 0:
            worked_dividends = True


    except:

        pass 

    try:
        # Titles
        titles1 = get_titles('Tableau_Histo_Valo')
        titles2 = get_titles('Tableau_Histo_ECR_a')
        titles4 = get_titles('Tableau_Histo_Sit_Fin')

        # Year
        years1 = get_years('Tableau_Histo_Valo')
        years2 = get_years('Tableau_Histo_ECR_a')
        years4 = get_years('Tableau_Histo_Sit_Fin')

        # Numbers
        numbers1 = get_numbers('Tableau_Histo_Valo')
        numbers2 = get_numbers('Tableau_Histo_ECR_a')
        numbers4 = get_numbers('Tableau_Histo_Sit_Fin')



        # Plot dataframes and values
        df1 = assemble_dataframe(title=titles1 , values=numbers1 , column=years1)
        dsi.export(df1 , 'data\\df1.png' , table_conversion='matplolib' , fontsize=10)
        if english == False:
            add_picture('data\\df1.png' , 'Tableau 1' , left=1.1 , top=2.1)
        else:
            add_picture('data\\df1.png' , 'Array 1' , left=1.1, top=2.1)
        plot_element(df1)

        df2 = assemble_dataframe(title=titles2 , values=numbers2 , column=years2)
        dsi.export(df2, 'data\\df2.png' , table_conversion='matplolib' , fontsize=12)
        if english == False:
            add_picture('data\\df2.png' , 'Tableau 2' ,  left=1, top=2.1)
        else:
            add_picture('data\\df2.png' , 'Array 2' , left=1, top=2.1)
        plot_element(df2)


        df4 = assemble_dataframe(title=titles4 , values=numbers4 , column=years4)

        dsi.export(df4 , 'data\\df4.png' , table_conversion='matplolib' , fontsize=11)
        if english == False:
            add_picture('data\\df4.png' , 'Tableau 3' , left=1.2, top=2.1)
        else:
            add_picture('data\\df4.png' , 'Array 3' , left=1.2, top=2.1)
        plot_element(df4)


        worked_dataframe = True

    except:
        pass


    print('Worked dataframes: ' , worked_dataframe)


    if worked_dataframe == True:
        
        try:
            payout_ratio(div=dividends , df2=df2)
        except :
            pass
        try:
            cap_vs_debt(df4)
        except :
            pass

        key_dataframe = get_key_df()
        dsi.export(key_dataframe , f'data\\key_dataframe.png' , table_conversion='matplolib')
        if english == False:
            add_picture('data\\key_dataframe.png' , 'Nombres clés' , left=2.2 , top=2)
        else:
            add_picture('data\\key_dataframe.png' , 'Key numbers' , left=2.2 , top=2)

    
    ###### Some visualisations ######

    if worked_stock_price == True:
        # Plot maximum draw down
        plot_maximum_draw_down(stock_price)
        # Plot against SP500
        plot_against_benmark(stock_price , sp500)
        # Plot a linear regression
        plot_regression(stock_price)
        try:
            # Plot the predictions
            NN_model(stock_price[['Volume',	'Adj Close']] , sp500[['Volume',	'Adj Close']] , vix)
        except TypeError:
            pass


        # Plot the relative strenght index
        plot_rsi(stock_price)
        # Plot the six months zoom
        plot_zoom_candles(stock_price)

    
    if worked_dividends == True and worked_stock_price == True:
        # Plot the price + dividends
        price_with_dividends(stock_price[['Adj Close']] , dividends)
    
    if worked_dividends == True:
         # Plot the historical dividends
        annual_dividend_history(dividends)



    print("Worked stock data : " , worked_stock_price)
    print("Worked dividends data : " , worked_dividends)

    try:
        # Plot the shareholders pie
        shareholders()
        worked_share_holders = True

    except:

        pass

    sentiment_score = False

    try:
        sentiment_scores()
        sentiment_score = True
    except:
        pass




    print('Sentiment score : ' , sentiment_score)
    print('Worked shareholders : ' , worked_share_holders)

    if worked_stock_price == True:

        ###### Somes interesting numbers ######
        stock_price = stock_price[['Adj Close']]
        ann_ , years = annualized_return(stock_price)
        ann_five , _years = annualized_return_five_years(stock_price)
        std_ = float(np.round((np.std(stock_price.pct_change())),4))
        new_index = stock_price.index[-1] - dt.timedelta(252 * 5)
        std_five = float(np.round((np.std(stock_price[new_index:].pct_change())),4))
        sharpe , sortino = sharpe_and_sortino_ratio(stock_price)
        sharpe_five , sortino_five = sharpe_and_sortino_ratio_five_years(stock_price)

        if worked_share_holders == True:
            try:
                main_inst = main_institutions()
            except:
                pass
            if english == False:
                try:
                    key_num = f'        Le retour sur investissement annuel sur {years} ans  est de {round(float(ann_) * 100 , 2)}%\n       et {round(float(ann_five) * 100 , 2)}% sur les 5 dernières années.\n        L\'écart type des gains et pertes est de {round(std_*100,4)}\n      et {round(std_five*100,4)} sur les 5 dernières années.\n        Les ratios Sharpe et Sortino sont respectivement {round(float(sharpe),4)} et {round(float(sortino),4)}\n        et de {round(float(sharpe_five),4)} et {round(float(sortino_five),4)} sur les 5 dernières années.\n\n       Les trois principaux investisseurs institutionnels sont , par ordre décroissant de détention:\n        {main_inst[0]}\n        {main_inst[1]}\n        {main_inst[2]}.\n\n\n\n\n       Date de ce rapport : {dt.datetime.today().date()}'
                except:
                    key_num = f'        Le retour sur investissement annuel sur {years} ans  est de {round(float(ann_) * 100 , 2)}%\n       et {round(float(ann_five) * 100 , 2)}% sur les 5 dernières années.\n        L\'écart type des gains et pertes est de {round(std_*100,4)}\n      et {round(std_five*100,4)} sur les 5 dernières années.\n        Les ratios Sharpe et Sortino sont respectivement {round(float(sharpe),4)} et {round(float(sortino),4)}\n        et de {round(float(sharpe_five),4)} et {round(float(sortino_five),4)} sur les 5 dernières années.\n\n\n\n\n       Date de ce rapport : {dt.datetime.today().date()}'
                
                description(key_num , 'Chiffres et informations à prendre en compte.')
            else:
                key_num = t(f'      Le retour sur investissement annuel sur {years} ans  est de {round(float(ann_) * 100 , 2)}%\n       et {round(float(ann_five) * 100 , 2)}% sur les 5 dernières années.\n        L\'écart type des gains et pertes est de {round(std_*100,4)}\n      et {round(std_five*100,4)} sur les 5 dernières années.\n        Les ratios Sharpe et Sortino sont respectivement {round(float(sharpe),4)} et {round(float(sortino),4)}\n        et de {round(float(sharpe_five),4)} et {round(float(sortino_five),4)} sur les 5 dernières années.\n\n       Les trois principaux investisseurs institutionnels sont , par ordre décroissant de détention:\n        {main_inst[0]}\n        {main_inst[1]}\n        {main_inst[2]}.\n\n\n\n\n        Date de ce rapport : {dt.datetime.today().date()}')
                description(key_num , ('Numbers and interestings informations.'))

        else:
            if english == False:
                key_num = key_num = f'      Le retour sur investissement annuel sur {years} ans  est de {round(float(ann_) * 100 , 2)}%\n        et {round(float(ann_five) * 100 , 2)}% sur les 5 dernières années.\n       L\'écart type des gains et pertes est de {round(std_*100,4)}\n      et {round(std_five*100,4)} sur les 5 dernières années.\n        Les ratios Sharpe et Sortino sont respectivement {round(float(sharpe),4)} et {round(float(sortino),4)}\n        et de {round(float(sharpe_five),4)} et {round(float(sortino_five),4)} sur les 5 dernières années.\n\n\n\n\n        Date de ce rapport : {dt.datetime.today().date()}'
                description(key_num , 'Chiffres et informations à prendre en compte.')
            else:
                key_num = t(f'      Le retour sur investissement annuel sur {years} ans  est de {round(float(ann_) * 100 , 2)}%\n        et {round(float(ann_five) * 100 , 2)}% sur les 5 dernières années.\n       L\'écart type des gains et pertes est de {round(std_*100,4)}\n      et {round(std_five*100,4)} sur les 5 dernières années.\n        Les ratios Sharpe et Sortino sont respectivement {round(float(sharpe),4)} et {round(float(sortino),4)}\n        et de {round(float(sharpe_five),4)} et {round(float(sortino_five),4)} sur les 5 dernières années.\n\n\n\n\n        Date de ce rapport : {dt.datetime.today().date()}')
                description(key_num , 'Numbers and interestings informations.')
        

    ##### Process the inputs #####
    ticker = ticker.upper()
    path = f'{path}\\{ticker}.pptx'.replace('"' , '')

    # Save the presentation
    pres.save(path)


############### GUI ###############
from tkinter import *
from tkinter import ttk

root = Tk()
root.title('Analyse')
label_language = Label(root , text='Language').grid(row=0 , column=0)
label_ticker = Label(root , text='Ticker de l\'entreprise / ticker of the company').grid(row=1 , column=0)
label_name = Label(root , text='Nom de l\'entreprise / name of the companie').grid(row=2 , column=0)
label_path = Label(root , text='Chemin de sauvegarde / save path').grid(row=3 , column=0)
label_button = Label(root , text='Débuter l\'analyse / Start the analysis').grid(row=4 , column=0)
language = ttk.Combobox(root, values=['Français' , 'English'] , width=22)
language.grid(row=0 , column=1)
ticker = Entry(root , width=25 , borderwidth=3)
ticker.grid(row=1 , column=1)
name = Entry(root , width=25, borderwidth=3)
name.grid(row=2 , column=1)
path = Entry(root , width=25, borderwidth=3)
path.grid(row=3 , column=1)


def get_entries(ticker=ticker , path=path):
    '''Triggered by the analysis button. Get the entries and start the run function.'''
    language_ = language.get()
    ticker_ = ticker.get()
    path_ = path.get()
    name_ = name.get()
    run(language=language_ , ticker=ticker_ , path=path_ , companie_name=name_)
    root.quit()

button = Button(root , text='Analyse / Analysis' , command=get_entries , width=21 , height=0 , borderwidth=1).grid(row=4 , column=1)
root.mainloop()